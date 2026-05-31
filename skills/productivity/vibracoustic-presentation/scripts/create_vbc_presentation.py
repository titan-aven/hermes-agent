#!/usr/bin/env python3
"""
Vibracoustic Presentation Builder — VALIDATED WORKING SCRIPT
============================================================
Basiert auf V2 (Kai_Voges_Orga20_Workshop_Jun16_v2.pptx) — die erste
korrekte VBC-Präsentation. Niemals dieses Script durch Neugenerierung
ersetzen — nur parametrieren und befüllen.

Grundregeln (eingebaut, nicht umgehbar):
  - Nur echte VC-Layouts aus dem Template
  - Nur Placeholders — keine freien Shapes/TextBoxes
  - Validator läuft automatisch vor dem Speichern

Usage:
  uv run --with python-pptx python3 create_vbc_presentation.py
"""

import sys
from pathlib import Path
from datetime import date
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

# ─── KONFIGURATION ────────────────────────────────────────────────────────────

TEMPLATE = "/tmp/vibracoustic-brand/VC_Brand_Kit.pptx"
OUTPUT   = "/tmp/vibracoustic-brand/output.pptx"

# Globale Metadaten — anpassen pro Präsentation
META = {
    "titel":    "Präsentationstitel",
    "topline":  "Bereich | Thema | Datum",       # Topline über dem Titel (Slide 2+)
    "untertitel": f"Kai Voges  |  {date.today().strftime('%d.%m.%Y')}",
    "datum":    date.today().strftime("%m/%d/%Y"),
}

# ─── INHALTE ─────────────────────────────────────────────────────────────────
# Jede Folie ist ein dict mit 'layout' und layout-spezifischen Keys.
# Verwende NUR die unten dokumentierten Layouts — niemals 'Blank' oder 'Only Title'.

SLIDES = [
    # ── Titelfolie ────────────────────────────────────────────────────────────
    {
        "layout": "Title Page",
        "titel": META["titel"],
        "untertitel": META["untertitel"],
    },

    # ── Standard-Inhaltsfolie (1 Content) ────────────────────────────────────
    # Für Text, Bullets, Tabellen — die häufigste Folie
    {
        "layout": "1 Content",
        "topline": META["topline"],
        "titel": "Agenda",
        "content": (
            "01  |  Erster Punkt\n"
            "02  |  Zweiter Punkt\n"
            "03  |  Dritter Punkt"
        ),
    },

    # ── 2×2 Raster (4 Contents with Headlinebox) ─────────────────────────────
    # Für 4 gleichwertige Kategorien/Themen
    {
        "layout": "4 Contents with Headlinebox",
        "topline": META["topline"],
        "titel": "Vier Bereiche",
        "col1_header": "Bereich A",   "col1": "Inhalt A",
        "col2_header": "Bereich B",   "col2": "Inhalt B",
        "col3_header": "Bereich C",   "col3": "Inhalt C",
        "col4_header": "Bereich D",   "col4": "Inhalt D",
    },

    # ── Zwei Spalten mit Überschriften (2 Contents with Headlinebox) ──────────
    # Für Vergleiche, Gegenüberstellungen
    {
        "layout": "2 Contents with Headlinebox",
        "topline": META["topline"],
        "titel": "Vergleich",
        "col1_header": "Spalte Links",
        "col1": "Inhalt links",
        "col2_header": "Spalte Rechts",
        "col2": "Inhalt rechts",
    },

    # ── 2×3 Raster (6 Contents with Headlinebox) ─────────────────────────────
    # Für 6 gleichwertige Punkte/Prinzipien
    {
        "layout": "6 Contents with Headlinebox",
        "topline": META["topline"],
        "titel": "Sechs Prinzipien",
        "col1_header": "01", "col1": "Erster Punkt",
        "col2_header": "02", "col2": "Zweiter Punkt",
        "col3_header": "03", "col3": "Dritter Punkt",
        "col4_header": "04", "col4": "Vierter Punkt",
        "col5_header": "05", "col5": "Fünfter Punkt",
        "col6_header": "06", "col6": "Sechster Punkt",
    },
]

# ─── TEMPLATE-LOADER ─────────────────────────────────────────────────────────

def load_template(path=TEMPLATE):
    """Lädt Template und löscht alle Beispiel-Slides."""
    if not Path(path).exists():
        print(f"Template nicht gefunden: {path}")
        print("Lade aus Drive...")
        import subprocess
        Path("/tmp/vibracoustic-brand").mkdir(exist_ok=True)
        subprocess.run([
            "gog", "drive", "download", "1PDGluRzlUhiOIhWqfA7UCrU2oIqwGkwX",
            "--account", "aven73.claw@gmail.com",
            "--output", path
        ], check=True)

    prs = Presentation(path)

    # Alle Template-Beispiel-Slides entfernen
    xml_slides = prs.slides._sldIdLst
    while len(xml_slides) > 0:
        rId = xml_slides[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        if rId:
            prs.part.drop_rel(rId)
        del xml_slides[0]

    return prs

def get_layouts(prs):
    """Gibt ein Dict {layout_name: layout} zurück."""
    return {
        l.name: l
        for master in prs.slide_masters
        for l in master.slide_layouts
    }

# ─── FOOTER HELPER ───────────────────────────────────────────────────────────

def set_footer(slide, titel=META["titel"], datum=META["datum"]):
    """Setzt Datum, Fußzeile und lässt Seitenzahl automatisch."""
    for ph in slide.placeholders:
        try:
            pt = ph.placeholder_format.type.name
            if pt == "DATE":
                ph.text = datum
            elif pt == "FOOTER":
                ph.text = f"|  {titel}"
        except Exception:
            pass

# ─── SLIDE BUILDER ───────────────────────────────────────────────────────────

def build_title_page(prs, layouts, data):
    slide = prs.slides.add_slide(layouts["Title Page"])
    slide.placeholders[0].text = data["titel"]
    slide.placeholders[1].text = data.get("untertitel", META["untertitel"])
    return slide


def build_1_content(prs, layouts, data):
    slide = prs.slides.add_slide(layouts["1 Content"])
    slide.placeholders[18].text = data.get("topline", META["topline"])
    slide.placeholders[0].text  = data["titel"]
    _fill_text(slide.placeholders[1], data.get("content", ""))
    set_footer(slide)
    return slide


def build_4_contents(prs, layouts, data):
    slide = prs.slides.add_slide(layouts["4 Contents with Headlinebox"])
    slide.placeholders[18].text = data.get("topline", META["topline"])
    slide.placeholders[0].text  = data["titel"]
    # Überschriften
    slide.placeholders[41].text = data.get("col1_header", "")
    slide.placeholders[28].text = data.get("col2_header", "")
    slide.placeholders[43].text = data.get("col3_header", "")
    slide.placeholders[42].text = data.get("col4_header", "")
    # Inhalte
    _fill_text(slide.placeholders[37], data.get("col1", ""))
    _fill_text(slide.placeholders[38], data.get("col2", ""))
    _fill_text(slide.placeholders[39], data.get("col3", ""))
    _fill_text(slide.placeholders[40], data.get("col4", ""))
    set_footer(slide)
    return slide


def build_2_contents_headline(prs, layouts, data):
    slide = prs.slides.add_slide(layouts["2 Contents with Headlinebox"])
    slide.placeholders[18].text = data.get("topline", META["topline"])
    slide.placeholders[0].text  = data["titel"]
    slide.placeholders[15].text = data.get("col1_header", "")
    slide.placeholders[28].text = data.get("col2_header", "")
    _fill_text(slide.placeholders[26], data.get("col1", ""))
    _fill_text(slide.placeholders[27], data.get("col2", ""))
    set_footer(slide)
    return slide


def build_6_contents(prs, layouts, data):
    slide = prs.slides.add_slide(layouts["6 Contents with Headlinebox"])
    slide.placeholders[18].text = data.get("topline", META["topline"])
    slide.placeholders[0].text  = data["titel"]
    # Überschriften Zeile 1
    slide.placeholders[35].text = data.get("col1_header", "")
    slide.placeholders[28].text = data.get("col2_header", "")
    slide.placeholders[56].text = data.get("col3_header", "")
    # Überschriften Zeile 2
    slide.placeholders[58].text = data.get("col4_header", "")
    slide.placeholders[57].text = data.get("col5_header", "")
    slide.placeholders[59].text = data.get("col6_header", "")
    # Inhalte
    _fill_text(slide.placeholders[47], data.get("col1", ""))
    _fill_text(slide.placeholders[48], data.get("col2", ""))
    _fill_text(slide.placeholders[49], data.get("col3", ""))
    _fill_text(slide.placeholders[53], data.get("col4", ""))
    _fill_text(slide.placeholders[54], data.get("col5", ""))
    _fill_text(slide.placeholders[55], data.get("col6", ""))
    set_footer(slide)
    return slide


def build_chapter_divider(prs, layouts, data):
    slide = prs.slides.add_slide(layouts["Chapter Divider"])
    slide.placeholders[0].text  = data["titel"]
    slide.placeholders[14].text = data.get("nummer", "01")
    set_footer(slide, titel=data.get("footer_titel", META["titel"]))
    return slide


def build_2_contents(prs, layouts, data):
    slide = prs.slides.add_slide(layouts["2 Contents"])
    slide.placeholders[21].text = data.get("topline", META["topline"])
    slide.placeholders[0].text  = data["titel"]
    _fill_text(slide.placeholders[19], data.get("col1", ""))
    _fill_text(slide.placeholders[20], data.get("col2", ""))
    set_footer(slide)
    return slide


# ─── TEXT HELPER ─────────────────────────────────────────────────────────────

def _fill_text(placeholder, text):
    """
    Füllt einen Placeholder mit Text.
    WICHTIG: Schreibt direkt in den Placeholder — niemals neue TextBoxes erstellen.
    Mehrzeilige Inhalte mit \\n trennen. Einrückung mit Level 1 = '  ' Prefix im Text.
    """
    if not text:
        return
    tf = placeholder.text_frame
    tf.word_wrap = True

    lines = text.split("\n")
    for j, line in enumerate(lines):
        # Level aus Einrückung bestimmen
        level = 0
        stripped = line.lstrip()
        if line.startswith("    "):
            level = 1
        if j == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text  = stripped
        para.level = level

# ─── LAYOUT DISPATCHER ───────────────────────────────────────────────────────

BUILDERS = {
    "Title Page":                    build_title_page,
    "1 Content":                     build_1_content,
    "4 Contents with Headlinebox":   build_4_contents,
    "2 Contents with Headlinebox":   build_2_contents_headline,
    "6 Contents with Headlinebox":   build_6_contents,
    "Chapter Divider":               build_chapter_divider,
    "2 Contents":                    build_2_contents,
}

# ─── VALIDATOR ───────────────────────────────────────────────────────────────

def validate(prs):
    """
    Prüft ob die Präsentation VBC-konform ist.
    Wirft ValueError bei Verstößen — Datei wird NICHT gespeichert.
    """
    errors = []
    BANNED_LAYOUTS = {"Blank", "Only Title"}

    for i, slide in enumerate(prs.slides, 1):
        layout = slide.slide_layout.name

        if layout in BANNED_LAYOUTS:
            errors.append(f"Slide {i}: Verbotenes Layout '{layout}'")

        free = [s for s in slide.shapes if s.shape_type != 14]
        if len(free) > 2:
            errors.append(
                f"Slide {i}: {len(free)} freie Shapes ({[s.name for s in free]}) "
                f"— nur Placeholders erlaubt"
            )

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name and run.font.name not in ("Tahoma", None):
                        errors.append(
                            f"Slide {i} / '{shape.name}': "
                            f"Schrift '{run.font.name}' — nur Tahoma erlaubt"
                        )

    if errors:
        print("\n⛔ VALIDATION FAILED — Datei wird nicht gespeichert:")
        for e in errors:
            print(f"   • {e}")
        raise ValueError(f"{len(errors)} Validierungsfehler")

    print(f"✅ Validation OK — {len(prs.slides)} Slides, alle VBC-konform")

# ─── MAIN ────────────────────────────────────────────────────────────────────

def build(slides_data=SLIDES, output=OUTPUT, meta=META):
    prs     = load_template()
    layouts = get_layouts(prs)

    for data in slides_data:
        layout_name = data["layout"]
        builder = BUILDERS.get(layout_name)
        if not builder:
            raise ValueError(
                f"Unbekanntes Layout '{layout_name}'. "
                f"Erlaubt: {list(BUILDERS.keys())}"
            )
        builder(prs, layouts, data)

    # Validator läuft immer — kann nicht übersprungen werden
    validate(prs)

    Path(output).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"💾 Gespeichert: {output}")
    return output


if __name__ == "__main__":
    build()
